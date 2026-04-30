"""
slide_capture.py — HTML 슬라이드를 Chrome으로 렌더링해 스크린샷 → PPTX 변환
웹에서 보이는 것과 100% 동일한 시각 품질을 보장한다.
"""

import time
import tempfile
from io import BytesIO
from pathlib import Path

# 슬라이드 캡처 해상도: 16:9, 고해상도
CAPTURE_W = 1920
CAPTURE_H = 1080


# ── 캡처용 CSS 주입 ─────────────────────────────────────────────────────
# 전환 애니메이션 제거, 숨겨진 요소 강제 표시
_CAPTURE_CSS = f"""
<style id="capture-override">
html, body {{
  width:  {CAPTURE_W}px !important;
  height: {CAPTURE_H}px !important;
  margin: 0 !important;
  padding: 0 !important;
  overflow: hidden !important;
  background: #0a0a0f !important;
}}
.slide-viewport, #viewport {{
  position: fixed !important;
  inset: 0 !important;
  width:  {CAPTURE_W}px !important;
  height: {CAPTURE_H}px !important;
  background: #0a0a0f !important;
  overflow: hidden !important;
}}
*, *::before, *::after {{
  transition: none !important;
  animation: none !important;
  animation-duration: 0s !important;
}}
/* reveal/appear 요소 즉시 표시 */
.reveal, [class*="reveal"], .appear, [data-reveal] {{
  opacity: 1 !important;
  transform: none !important;
  filter: none !important;
  visibility: visible !important;
}}
</style>
"""


def _inject_capture_css(html: str) -> str:
    return html.replace('</head>', _CAPTURE_CSS + '</head>', 1)


# ── JavaScript: 슬라이드 즉시 전환 ────────────────────────────────────
def _go_slide_js(i: int) -> str:
    """CSS transition 없이 index i 슬라이드를 즉시 표시하는 JS"""
    return f"""
        (function() {{
            const vp = document.getElementById('viewport') ||
                       document.querySelector('.slide-viewport');
            if (!vp) return;

            // 뷰포트 강제 풀스크린
            vp.style.cssText += '; position:fixed !important; inset:0 !important; ' +
                                 'width:{CAPTURE_W}px !important; height:{CAPTURE_H}px !important;';

            const slides = vp.querySelectorAll('.slide');
            const idx = {i};

            slides.forEach((s, n) => {{
                s.style.transition = 'none';
                s.classList.remove('active', 'stand-left');
                if (n < idx) {{
                    s.style.cssText += '; opacity:0; transform:translateX(-100%); pointer-events:none;';
                    s.classList.add('stand-left');
                }} else if (n === idx) {{
                    s.style.cssText += '; opacity:1; transform:translateX(0); pointer-events:auto;';
                    s.classList.add('active');
                    // 해당 슬라이드의 모든 reveal 요소 즉시 표시
                    s.querySelectorAll('.reveal, [class*="reveal"], .appear').forEach(el => {{
                        el.style.cssText += '; opacity:1 !important; transform:none !important; filter:none !important;';
                    }});
                }} else {{
                    s.style.cssText += '; opacity:0; transform:translateX(100%); pointer-events:none;';
                }}
            }});
        }})();
    """


def capture_slides(html: str, progress_callback=None) -> list[bytes]:
    """
    HTML 프레젠테이션 → 슬라이드별 PNG 스크린샷 리스트.

    Args:
        html: 완성된 HTML 프레젠테이션 문자열
        progress_callback: (current, total) 형태로 호출되는 진행 콜백 (선택)

    Returns:
        PNG 바이트 리스트 (슬라이드 순서대로, 각 1920×1080)
    """
    try:
        from selenium import webdriver
        from selenium.webdriver.chrome.options import Options
        from selenium.webdriver.chrome.service import Service
        from webdriver_manager.chrome import ChromeDriverManager
    except ImportError as e:
        raise RuntimeError(
            f"selenium 또는 webdriver-manager 미설치: {e}\n"
            "pip install selenium webdriver-manager 실행 후 재시도하세요."
        ) from e

    capture_html = _inject_capture_css(html)

    tmp = tempfile.NamedTemporaryFile(
        mode='w', suffix='.html', delete=False, encoding='utf-8'
    )
    tmp.write(capture_html)
    tmp.close()
    tmp_path = Path(tmp.name)

    options = Options()
    options.add_argument('--headless=new')
    options.add_argument('--no-sandbox')
    options.add_argument('--disable-dev-shm-usage')
    options.add_argument('--disable-gpu')
    options.add_argument(f'--window-size={CAPTURE_W},{CAPTURE_H}')
    options.add_argument('--hide-scrollbars')
    options.add_argument('--force-device-scale-factor=1')

    driver = None
    try:
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=options)

        # CDP로 정확한 뷰포트 크기 강제 지정
        driver.execute_cdp_cmd('Emulation.setDeviceMetricsOverride', {
            'width':             CAPTURE_W,
            'height':            CAPTURE_H,
            'deviceScaleFactor': 1,
            'mobile':            False,
        })

        file_url = 'file:///' + tmp_path.as_posix().lstrip('/')
        driver.get(file_url)

        # 폰트·이미지 로딩 대기
        time.sleep(2.5)

        # html/body가 정확히 1920×1080을 채우도록 JS로 한 번 더 강제
        driver.execute_script(f"""
            document.documentElement.style.cssText =
                'margin:0;padding:0;overflow:hidden;width:{CAPTURE_W}px;height:{CAPTURE_H}px;';
            document.body.style.cssText =
                'margin:0;padding:0;overflow:hidden;width:{CAPTURE_W}px;height:{CAPTURE_H}px;background:#0a0a0f;';
            var vp = document.getElementById('viewport') ||
                     document.querySelector('.slide-viewport');
            if (vp) {{
                vp.style.cssText =
                    'position:fixed;top:0;left:0;width:{CAPTURE_W}px;height:{CAPTURE_H}px;' +
                    'overflow:hidden;z-index:1;background:#0a0a0f;';
            }}
        """)

        # 총 슬라이드 수
        total = driver.execute_script(
            'return (document.getElementById("viewport") || document.querySelector(".slide-viewport") || document.body)'
            '.querySelectorAll(".slide").length'
        )
        if not total:
            raise RuntimeError("슬라이드를 찾을 수 없습니다. HTML 구조를 확인하세요.")

        screenshots: list[bytes] = []
        for i in range(total):
            if progress_callback:
                progress_callback(i, total)

            # CSS transition 없이 슬라이드 즉시 전환
            driver.execute_script(_go_slide_js(i))
            time.sleep(0.3)

            # CDP가 창을 1920×1080으로 고정했으므로 전체 화면 스크린샷 = 정확히 1920×1080
            # (요소 단위 캡처는 요소 크기에 따라 작게 찍혀 빈 여백이 생기므로 사용하지 않음)
            screenshots.append(_ensure_size(driver.get_screenshot_as_png()))

        return screenshots

    finally:
        if driver:
            try:
                driver.quit()
            except Exception:
                pass
        tmp_path.unlink(missing_ok=True)


def _ensure_size(png_bytes: bytes) -> bytes:
    """PNG를 정확히 CAPTURE_W × CAPTURE_H 로 리사이즈 (Pillow 있을 때만)"""
    try:
        from PIL import Image
        img = Image.open(BytesIO(png_bytes))
        if img.size != (CAPTURE_W, CAPTURE_H):
            img = img.resize((CAPTURE_W, CAPTURE_H), Image.LANCZOS)
            buf = BytesIO()
            img.save(buf, format='PNG')
            return buf.getvalue()
    except ImportError:
        pass  # Pillow 없으면 그냥 원본 반환
    return png_bytes


def screenshots_to_pptx(screenshots: list[bytes]) -> bytes:
    """PNG 스크린샷 리스트 → PPTX 바이트 (슬라이드 전체를 이미지로 채움)"""
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.dml.color import RGBColor

    # 16:9 와이드 슬라이드 (1920:1080 비율과 동일)
    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    for png_bytes in screenshots:
        slide = prs.slides.add_slide(blank_layout)

        # 배경 검정 (이미지 사이 흰 틈 방지)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0F)

        # 슬라이드 전체를 스크린샷으로 채우기 (0,0 부터 슬라이드 전체 크기)
        slide.shapes.add_picture(
            BytesIO(png_bytes),
            Emu(0), Emu(0),
            SLIDE_W, SLIDE_H,
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def html_to_pptx(html: str, progress_callback=None) -> bytes:
    """HTML 프레젠테이션 → PPTX (스크린샷 기반, 웹과 100% 시각 동일)"""
    shots = capture_slides(html, progress_callback=progress_callback)
    return screenshots_to_pptx(shots)
