"""
ppt_generator.py — 강의용 PPT 생성기
Claude가 정리한 슬라이드 데이터를 받아 .pptx 파일로 변환한다.
"""

import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── 디자인 상수 ────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0F, 0x17, 0x2A)   # 메인 배경 (진한 남색)
BG_SECTION  = RGBColor(0xF9, 0x73, 0x16)   # 섹션 구분 배경 (오렌지)
BG_HIGHL    = RGBColor(0x1E, 0x29, 0x3B)   # 강조 슬라이드 배경 (중간 남색)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xF9, 0x73, 0x16)
LIGHT_GRAY  = RGBColor(0xCB, 0xD5, 0xE1)
YELLOW      = RGBColor(0xFB, 0xBF, 0x24)
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)


# ── 유틸 ──────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(slide, text, left, top, w, h,
              size=24, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_rect(slide, left, top, w, h, fill_color: RGBColor, radius=False):
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _strip_md(text: str) -> str:
    """마크다운 기호 제거"""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text.strip()


def _bold_parts(text: str):
    """**굵게** 부분을 (text, bold) 튜플 리스트로 분리"""
    parts = re.split(r"(\*\*.+?\*\*)", text)
    result = []
    for p in parts:
        if p.startswith("**") and p.endswith("**"):
            result.append((p[2:-2], True))
        elif p:
            result.append((p, False))
    return result


def _add_text_with_bold(slide, text, left, top, w, h,
                        size=24, base_color=WHITE, bold_color=YELLOW,
                        align=PP_ALIGN.LEFT, wrap=True):
    """**강조** 가 있을 경우 볼드+색상으로 렌더링"""
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    for part_text, is_bold in _bold_parts(text):
        run = p.add_run()
        run.text = part_text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = bold_color if is_bold else base_color
    return box


# ── 슬라이드 유형별 생성 ───────────────────────────────────────────────

def _slide_title(prs, title: str, subtitle: str = ""):
    """표지 슬라이드"""
    slide = _blank_slide(prs)
    _set_bg(slide, BG_DARK)

    # 왼쪽 오렌지 바
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, ORANGE)

    # 채널명
    _add_text(slide, "로직해커 엑스", Inches(0.5), Inches(1.6),
              Inches(12), Inches(0.6), size=18, color=ORANGE, bold=True)

    # 메인 제목
    display = title.split("|")[0].strip() if "|" in title else title
    _add_text(slide, display, Inches(0.5), Inches(2.3),
              Inches(12), Inches(2.8), size=38, bold=True, color=WHITE, wrap=True)

    # 구분선
    _add_rect(slide, Inches(0.5), Inches(5.3), Inches(4), Pt(3), ORANGE)

    if subtitle:
        _add_text(slide, subtitle, Inches(0.5), Inches(5.6),
                  Inches(12), Inches(0.8), size=20, color=LIGHT_GRAY)


def _slide_section(prs, section_name: str, slide_title: str):
    """섹션 구분 슬라이드 — 오렌지 배경"""
    slide = _blank_slide(prs)
    _set_bg(slide, BG_SECTION)

    # 섹션 번호 느낌의 큰 배경 텍스트
    _add_text(slide, section_name, Inches(0.5), Inches(2.0),
              Inches(12), Inches(1.0), size=22, color=WHITE, bold=False,
              align=PP_ALIGN.CENTER)

    # 구분선
    _add_rect(slide, Inches(4.5), Inches(3.2), Inches(4.3), Pt(3),
              RGBColor(0xFF, 0xFF, 0xFF))

    _add_text(slide, slide_title, Inches(0.5), Inches(3.5),
              Inches(12), Inches(1.8), size=40, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, wrap=True)


def _slide_bullets(prs, title: str, points: list):
    """핵심 포인트 슬라이드 — 불릿 리스트"""
    slide = _blank_slide(prs)
    _set_bg(slide, BG_DARK)

    # 상단 오렌지 바
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ORANGE)

    # 슬라이드 제목
    _add_text_with_bold(slide, title, Inches(0.6), Inches(0.3),
                        Inches(11.5), Inches(1.0),
                        size=30, base_color=WHITE, bold_color=YELLOW,
                        align=PP_ALIGN.LEFT)

    # 구분선
    _add_rect(slide, Inches(0.6), Inches(1.4), Inches(11), Pt(2), ORANGE)

    # 불릿 포인트들
    y_positions = [Inches(1.7), Inches(2.65), Inches(3.6), Inches(4.55)]
    for i, point in enumerate(points[:4]):
        y = y_positions[i]
        # 오렌지 동그라미 마커
        _add_rect(slide, Inches(0.55), y + Inches(0.12),
                  Pt(14), Pt(14), ORANGE)
        # 포인트 텍스트
        _add_text_with_bold(slide, point.lstrip("•- ").strip(),
                            Inches(0.95), y, Inches(11.5), Inches(0.85),
                            size=26, base_color=LIGHT_GRAY, bold_color=YELLOW)


def _slide_highlight(prs, quote: str, desc: str = ""):
    """임팩트 강조 슬라이드 — 한 줄 핵심 문구"""
    slide = _blank_slide(prs)
    _set_bg(slide, BG_HIGHL)

    # 큰 따옴표 장식
    _add_text(slide, "❝", Inches(0.5), Inches(0.8),
              Inches(2), Inches(1.5), size=72, color=ORANGE, bold=True)

    # 핵심 문구
    _add_text_with_bold(slide, quote, Inches(0.6), Inches(2.0),
                        Inches(12), Inches(2.2),
                        size=36, base_color=WHITE, bold_color=YELLOW,
                        align=PP_ALIGN.CENTER, wrap=True)

    if desc:
        _add_text(slide, _strip_md(desc), Inches(0.6), Inches(5.0),
                  Inches(12), Inches(1.2), size=20, color=LIGHT_GRAY,
                  align=PP_ALIGN.CENTER)


# ── 슬라이드 데이터 파서 ──────────────────────────────────────────────

def _parse_slide_data(raw: str) -> list:
    """
    Claude가 반환한 텍스트를 슬라이드 딕셔너리 리스트로 파싱한다.
    """
    slides = []
    blocks = re.split(r"\[(?:TITLE|SECTION|BULLETS|HIGHLIGHT)\]", raw)
    tags   = re.findall(r"\[(TITLE|SECTION|BULLETS|HIGHLIGHT)\]", raw)

    for tag, block in zip(tags, blocks[1:]):
        lines = [l.strip() for l in block.strip().splitlines() if l.strip()]
        data  = {"type": tag.lower()}

        for line in lines:
            if ":" in line:
                key, _, val = line.partition(":")
                data[key.strip()] = val.strip()
            elif line.startswith("•") or line.startswith("-"):
                data.setdefault("points", []).append(line)

        slides.append(data)

    return slides


# ── 메인 함수 ─────────────────────────────────────────────────────────

def generate_ppt(title: str, slide_content: str) -> bytes:
    """
    title        : 강의 제목
    slide_content: Claude call_ppt_content() 반환값
    return       : .pptx bytes
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = _parse_slide_data(slide_content)

    for s in slides_data:
        t = s.get("type", "")

        if t == "title":
            _slide_title(prs,
                         s.get("제목", title),
                         s.get("부제", ""))

        elif t == "section":
            _slide_section(prs,
                           s.get("섹션명", ""),
                           s.get("슬라이드제목", ""))

        elif t == "bullets":
            points = s.get("points", [])
            _slide_bullets(prs,
                           s.get("슬라이드제목", ""),
                           points)

        elif t == "highlight":
            _slide_highlight(prs,
                             s.get("강조문구", ""),
                             s.get("설명", ""))

    # 슬라이드가 하나도 없으면 기본 표지만
    if len(prs.slides) == 0:
        _slide_title(prs, title)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
