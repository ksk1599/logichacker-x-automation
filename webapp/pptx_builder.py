"""
pptx_builder.py — HTML 슬라이드 → .pptx 변환
HTML의 <section> 요소를 파싱해 다크 테마 PowerPoint 파일을 생성한다.
"""

import base64
import re
from io import BytesIO

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── 색상 (template.html 디자인 시스템과 동일) ──────────────────────────
BG          = RGBColor(0x0A, 0x0A, 0x0F)
ACCENT      = RGBColor(0xC6, 0xA5, 0x5C)
TEXT_PRI    = RGBColor(0xF0, 0xEB, 0xE0)
TEXT_SEC    = RGBColor(0xB8, 0xB0, 0xA0)
TEXT_DIM    = RGBColor(0x70, 0x68, 0x58)

# ── 슬라이드 크기: 16:9 와이드 ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

MARGIN_L  = Inches(0.8)
MARGIN_T  = Inches(0.6)
MARGIN_R  = Inches(0.8)
CONTENT_W = W - MARGIN_L - MARGIN_R

# PPTX 네임스페이스
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# 사용할 한국어 폰트 (Windows 기본 내장)
KO_FONT = 'Malgun Gothic'


# ═══════════════════════════════════════════════════════════════════════
# 내부 유틸
# ═══════════════════════════════════════════════════════════════════════

def _strip_tags(html: str) -> str:
    """HTML 태그 제거 후 공백 정리"""
    text = re.sub(r'<[^>]+>', '', html)
    text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>') \
               .replace('&nbsp;', ' ').replace('&#39;', "'").replace('&quot;', '"')
    return re.sub(r'\s+', ' ', text).strip()


def _set_bg(slide, color: RGBColor):
    """슬라이드 배경 단색 채우기"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_image_from_b64(slide, data_url: str, left, top, width, height):
    """base64 data URL 이미지를 슬라이드에 추가"""
    try:
        b64_part = data_url.split(",", 1)[1]
        img_bytes = base64.b64decode(b64_part)
        slide.shapes.add_picture(BytesIO(img_bytes), left, top, width, height)
    except Exception:
        pass


def _add_rect(slide, left, top, width, height, color: RGBColor):
    """단색 직사각형 도형 추가"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def _add_text(slide, text, left, top, width, height,
              size=24, color=TEXT_PRI, bold=False,
              align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """텍스트박스 추가 후 스타일 적용. shape 반환."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = KO_FONT
    return txBox


def _add_bullet_list(slide, items: list[str], left, top, width, height,
                     size=20, color=TEXT_PRI):
    """불릿 리스트 텍스트박스 추가. shape 반환."""
    if not items:
        return None
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        dot = p.add_run()
        dot.text = '▸  '
        dot.font.size = Pt(size - 2)
        dot.font.color.rgb = ACCENT
        dot.font.name = KO_FONT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = KO_FONT

        # 항목 간 위 여백
        pPr = p._p.get_or_add_pPr()
        existing_spc = pPr.find(qn('a:spcBef'))
        if existing_spc is not None:
            pPr.remove(existing_spc)
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spc = etree.SubElement(spcBef, qn('a:spcPts'))
        spc.set('val', '160')

    return txBox


# ═══════════════════════════════════════════════════════════════════════
# 슬라이드 효과 (전환 & 애니메이션)
# ═══════════════════════════════════════════════════════════════════════

def _add_slide_transition(ppt_slide):
    """슬라이드 페이드 전환 효과 추가"""
    sld = ppt_slide._element
    trans = etree.SubElement(sld, f'{{{_P}}}transition', spd='med')
    etree.SubElement(trans, f'{{{_P}}}fade')


def _add_click_animations(ppt_slide, shape_ids: list[int]):
    """
    지정된 shape들이 클릭 순서대로 나타나는 애니메이션 추가.
    PowerPoint 'Appear' 효과(presetID=1) — 클릭 전 숨김, 클릭 시 표시.
    """
    if not shape_ids:
        return

    sld = ppt_slide._element
    _c = [1]

    def nid():
        v = _c[0]; _c[0] += 1; return str(v)

    def sub(parent, tag, **kw):
        return etree.SubElement(parent, f'{{{_P}}}{tag}',
                                 **{k: str(v) for k, v in kw.items()})

    timing  = sub(sld, 'timing')
    tnLst   = sub(timing, 'tnLst')
    par0    = sub(tnLst, 'par')
    cTn0    = sub(par0, 'cTn', id=nid(), dur='indefinite',
                   restart='whenNotActive', nodeType='tmRoot')
    cTnL0   = sub(cTn0, 'childTnLst')
    seq     = sub(cTnL0, 'seq', concurrent='1', nextAc='seek')
    cTn1    = sub(seq, 'cTn', id=nid(), dur='indefinite', nodeType='mainSeq')
    cTnL1   = sub(cTn1, 'childTnLst')

    for i, spid in enumerate(shape_ids):
        par      = sub(cTnL1, 'par')
        cTn_cl   = sub(par, 'cTn', id=nid(), fill='hold')
        stCL     = sub(cTn_cl, 'stCondLst')
        cond     = sub(stCL, 'cond', evt='onClick', delay='0')
        sub(cond, 'tn')

        cTnL2   = sub(cTn_cl, 'childTnLst')
        par2    = sub(cTnL2, 'par')
        cTn_fx  = sub(par2, 'cTn', id=nid(),
                       presetID='1', presetClass='entr', presetSubtype='0',
                       fill='hold', grpId=str(i), nodeType='clickEffect')
        stCL2   = sub(cTn_fx, 'stCondLst')
        sub(stCL2, 'cond', delay='0')

        cTnL3   = sub(cTn_fx, 'childTnLst')
        set_el  = sub(cTnL3, 'set')
        cBhvr   = sub(set_el, 'cBhvr')
        cTn_bh  = sub(cBhvr, 'cTn', id=nid(), dur='1', fill='hold')
        stCL3   = sub(cTn_bh, 'stCondLst')
        sub(stCL3, 'cond', delay='0')
        tgtEl   = sub(cBhvr, 'tgtEl')
        sub(tgtEl, 'spTgt', spid=str(spid))
        attrNL  = sub(cBhvr, 'attrNameLst')
        aName   = sub(attrNL, 'attrName')
        aName.text = 'style.visibility'
        to_el   = sub(set_el, 'to')
        sub(to_el, 'strVal', val='visible')

    # 이전/다음 클릭 조건 (슬라이드 탐색용)
    prevCL = sub(seq, 'prevCondLst')
    prev   = sub(prevCL, 'cond', evt='onPrevClick', delay='0')
    sub(prev, 'tn')
    nextCL = sub(seq, 'nextCondLst')
    nxt    = sub(nextCL, 'cond', evt='onNextClick', delay='0')
    sub(nxt, 'tn')

    sub(timing, 'bldLst')


# ═══════════════════════════════════════════════════════════════════════
# HTML 파싱
# ═══════════════════════════════════════════════════════════════════════

def _parse_slides(html: str) -> list[dict]:
    """<section class="slide ..."> 요소를 슬라이드 데이터로 파싱"""
    sections = re.findall(
        r'<section[^>]+class="([^"]*slide[^"]*)"[^>]*>(.*?)</section>',
        html, re.DOTALL | re.IGNORECASE
    )

    slides = []
    for classes, inner in sections:
        if 'slide--intro' in classes:
            stype = 'intro'
        elif 'slide--end' in classes:
            stype = 'end'
        elif 'slide--diagram' in classes:
            stype = 'diagram'
        elif 'slide--code' in classes:
            stype = 'code'
        else:
            stype = 'content'

        h1 = re.search(r'<h1[^>]*>(.*?)</h1>', inner, re.DOTALL)
        h2 = re.search(r'<h2[^>]*>(.*?)</h2>', inner, re.DOTALL)
        h3 = re.search(r'<h3[^>]*>(.*?)</h3>', inner, re.DOTALL)
        title = _strip_tags(
            (h1 or h2 or h3).group(1) if (h1 or h2 or h3) else ''
        )

        badge = re.search(
            r'class="terminal-badge[^"]*"[^>]*>(.*?)</div>', inner, re.DOTALL
        )
        badge_text = _strip_tags(badge.group(1)) if badge else ''

        items = [_strip_tags(m) for m in re.findall(r'<li[^>]*>(.*?)</li>', inner, re.DOTALL)]
        items = [x for x in items if x]

        no_list = re.sub(r'<ul[^>]*>.*?</ul>', '', inner, flags=re.DOTALL)
        no_list = re.sub(r'<ol[^>]*>.*?</ol>', '', no_list, flags=re.DOTALL)
        paras = [
            _strip_tags(m)
            for m in re.findall(r'<p[^>]*>(.*?)</p>', no_list, re.DOTALL)
        ]
        paras = [x for x in paras if x and len(x) > 1]

        cards = [
            _strip_tags(m)
            for m in re.findall(r'class="split-card[^"]*"[^>]*>(.*?)</div>', inner, re.DOTALL)
        ]
        cards = [x for x in cards if x]

        code_block = re.search(r'<code[^>]*>(.*?)</code>', inner, re.DOTALL)
        code_text = _strip_tags(code_block.group(1)) if code_block else ''

        img_urls = re.findall(r'<img[^>]+src="(data:[^"]+)"', inner, re.DOTALL)

        slides.append({
            'type':   stype,
            'title':  title,
            'badge':  badge_text,
            'items':  items,
            'paras':  paras,
            'cards':  cards,
            'code':   code_text,
            'images': img_urls,
        })

    return slides


# ═══════════════════════════════════════════════════════════════════════
# 슬라이드 타입별 렌더러 — 애니메이션할 shape_id 목록 반환
# ═══════════════════════════════════════════════════════════════════════

def _slide_intro(slide, data: dict) -> list[int]:
    """타이틀 슬라이드 (slide--intro)"""
    _set_bg(slide, BG)
    _add_rect(slide, Inches(0.4), Inches(1.5), Inches(0.06), Inches(4.5), ACCENT)

    if data['badge']:
        _add_text(slide, f'> {data["badge"]}',
                  MARGIN_L, Inches(1.6), CONTENT_W, Inches(0.5),
                  size=14, color=ACCENT, bold=True)

    _add_text(slide, data['title'],
              MARGIN_L, Inches(2.2), CONTENT_W, Inches(2.0),
              size=40, color=TEXT_PRI, bold=True)

    _add_rect(slide, MARGIN_L, Inches(4.4), Inches(1.2), Inches(0.04), ACCENT)

    if data['paras']:
        _add_text(slide, data['paras'][0],
                  MARGIN_L, Inches(4.6), CONTENT_W, Inches(0.8),
                  size=18, color=TEXT_SEC)

    _add_text(slide, '← → 키 또는 버튼으로 이동',
              MARGIN_L, Inches(6.6), CONTENT_W, Inches(0.5),
              size=12, color=TEXT_DIM)
    return []


def _slide_content(slide, data: dict) -> list[int]:
    """일반 콘텐츠 슬라이드. 클릭 애니메이션할 shape_id 목록 반환."""
    _set_bg(slide, BG)
    _add_rect(slide, 0, 0, W, Inches(0.04), ACCENT)

    anim_ids: list[int] = []
    top = MARGIN_T

    if data['badge']:
        _add_text(slide, f'> {data["badge"]}',
                  MARGIN_L, top, CONTENT_W, Inches(0.45),
                  size=13, color=ACCENT, bold=True)
        top += Inches(0.5)

    if data['title']:
        _add_text(slide, data['title'],
                  MARGIN_L, top, CONTENT_W, Inches(1.0),
                  size=28, color=TEXT_PRI, bold=True)
        top += Inches(1.0)

    _add_rect(slide, MARGIN_L, top, Inches(0.8), Inches(0.035), ACCENT)
    top += Inches(0.25)

    images = data.get('images', [])

    if images:
        half_w   = (CONTENT_W - Inches(0.4)) / 2
        text_top = top  # 텍스트 위치는 독립적으로 관리 (이미지와 분리)

        if data['items']:
            item_h = min(
                Inches(0.48) * len(data['items']) + Inches(0.4),
                H - text_top - Inches(0.5)
            )
            bx = _add_bullet_list(slide, data['items'],
                                   MARGIN_L, text_top, half_w, item_h, size=17)
            if bx:
                anim_ids.append(bx.shape_id)
            text_top += item_h

        for para in data['paras']:
            if text_top + Inches(0.7) > H - Inches(0.3):
                break
            bx = _add_text(slide, para,
                            MARGIN_L, text_top, half_w, Inches(0.7),
                            size=15, color=TEXT_SEC)
            anim_ids.append(bx.shape_id)
            text_top += Inches(0.65)

        # 오른쪽 이미지 (첫 번째만 사용)
        img_left = MARGIN_L + half_w + Inches(0.4)
        img_top  = MARGIN_T + Inches(1.4)
        img_h    = H - img_top - Inches(0.6)
        _add_image_from_b64(slide, images[0], img_left, img_top, half_w, img_h)

    else:
        if data['items']:
            item_h = min(
                Inches(0.5) * len(data['items']) + Inches(0.5),
                H - top - Inches(0.5)
            )
            bx = _add_bullet_list(slide, data['items'],
                                   MARGIN_L, top, CONTENT_W, item_h, size=20)
            if bx:
                anim_ids.append(bx.shape_id)
            top += item_h

        for para in data['paras']:
            if top + Inches(0.8) > H - Inches(0.3):
                break
            bx = _add_text(slide, para,
                            MARGIN_L, top, CONTENT_W, Inches(0.8),
                            size=18, color=TEXT_SEC)
            anim_ids.append(bx.shape_id)
            top += Inches(0.75)

    return anim_ids


def _slide_diagram(slide, data: dict) -> list[int]:
    """분할 레이아웃 슬라이드 (slide--diagram)"""
    _set_bg(slide, BG)
    _add_rect(slide, 0, 0, W, Inches(0.04), ACCENT)

    top = MARGIN_T
    if data['badge']:
        _add_text(slide, f'> {data["badge"]}',
                  MARGIN_L, top, CONTENT_W, Inches(0.45),
                  size=13, color=ACCENT, bold=True)
        top += Inches(0.5)

    if data['title']:
        _add_text(slide, data['title'],
                  MARGIN_L, top, CONTENT_W, Inches(0.9),
                  size=26, color=TEXT_PRI, bold=True)
        top += Inches(1.0)

    _add_rect(slide, MARGIN_L, top, Inches(0.8), Inches(0.035), ACCENT)
    top += Inches(0.3)

    card_w = (CONTENT_W - Inches(0.4)) / 2
    card_h = Inches(3.0)
    gap    = Inches(0.4)

    for i, card_text in enumerate(data['cards'][:2]):
        cx = MARGIN_L + i * (card_w + gap)
        bg_shape = slide.shapes.add_shape(1, cx, top, card_w, card_h)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x18)
        bg_shape.line.color.rgb = RGBColor(0x1E, 0x1C, 0x18)
        bg_shape.line.width = Pt(1)
        _add_text(slide, card_text,
                  cx + Inches(0.2), top + Inches(0.3),
                  card_w - Inches(0.4), card_h - Inches(0.4),
                  size=17, color=TEXT_PRI, wrap=True)
    return []


def _slide_code(slide, data: dict) -> list[int]:
    """코드 슬라이드 (slide--code)"""
    _set_bg(slide, BG)
    _add_rect(slide, 0, 0, W, Inches(0.04), ACCENT)

    top = MARGIN_T
    if data['title']:
        _add_text(slide, data['title'],
                  MARGIN_L, top, CONTENT_W, Inches(0.9),
                  size=26, color=ACCENT, bold=True)
        top += Inches(1.0)

    if data['code']:
        code_h = Inches(4.5)
        bg = slide.shapes.add_shape(1, MARGIN_L, top, CONTENT_W, code_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x18)
        bg.line.color.rgb = RGBColor(0x1E, 0x1C, 0x18)
        bg.line.width = Pt(1)
        _add_text(slide, data['code'],
                  MARGIN_L + Inches(0.2), top + Inches(0.2),
                  CONTENT_W - Inches(0.4), code_h - Inches(0.4),
                  size=15, color=TEXT_PRI, wrap=True)
    return []


def _slide_end(slide, data: dict) -> list[int]:
    """마무리 슬라이드 (slide--end)"""
    _set_bg(slide, BG)

    cx = W / 2
    _add_text(slide, data['title'] or '감사합니다',
              Inches(1), Inches(2.5), W - Inches(2), Inches(1.5),
              size=44, color=TEXT_PRI, bold=True, align=PP_ALIGN.CENTER)

    _add_rect(slide, cx - Inches(0.6), Inches(4.1), Inches(1.2), Inches(0.04), ACCENT)

    if data['paras']:
        _add_text(slide, data['paras'][0],
                  Inches(1), Inches(4.3), W - Inches(2), Inches(0.8),
                  size=20, color=TEXT_SEC, align=PP_ALIGN.CENTER)
    return []


# ═══════════════════════════════════════════════════════════════════════
# 공개 API
# ═══════════════════════════════════════════════════════════════════════

def build_pptx(html: str) -> bytes:
    """
    HTML 프레젠테이션 문자열 → .pptx 바이트 반환.
    앱에서 st.download_button(data=...) 에 바로 넣을 수 있다.
    """
    slides_data = _parse_slides(html)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]

    renderer = {
        'intro':   _slide_intro,
        'content': _slide_content,
        'diagram': _slide_diagram,
        'code':    _slide_code,
        'end':     _slide_end,
    }

    for data in slides_data:
        ppt_slide = prs.slides.add_slide(blank_layout)
        fn = renderer.get(data['type'], _slide_content)
        anim_ids = fn(ppt_slide, data)

        # 슬라이드 전환 효과 (페이드) — transition은 timing보다 먼저 추가
        _add_slide_transition(ppt_slide)

        # 클릭시 나타나기 애니메이션
        if anim_ids:
            _add_click_animations(ppt_slide, anim_ids)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
